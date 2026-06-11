# -*- coding: utf-8 -*-
"""Build the EduScheduler demonstration deck (non-technical audience)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1E, 0x27, 0x61)   # primary
INK    = RGBColor(0x14, 0x18, 0x33)   # darkest (title/closing bg)
ICE    = RGBColor(0xCA, 0xDC, 0xFC)   # light blue
AMBER  = RGBColor(0xF4, 0xA9, 0x3C)   # accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SLATE  = RGBColor(0x55, 0x5F, 0x86)   # muted text
INKTX  = RGBColor(0x1A, 0x1F, 0x36)   # body text on light
LIGHT  = RGBColor(0xEE, 0xF1, 0xFA)   # light content bg
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
TINT   = RGBColor(0xF6, 0xF8, 0xFD)   # subtle card tint
GREEN  = RGBColor(0x16, 0xA3, 0x4A)

HEAD = "Cambria"
BODY = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def _soft_shadow(shape, color=NAVY, alpha=86):
    sp = shape._element.spPr
    for tag in ('a:effectLst',):
        ex = sp.find(qn(tag))
        if ex is not None:
            sp.remove(ex)
    eff = sp.makeelement(qn('a:effectLst'), {})
    sh = eff.makeelement(qn('a:outerShdw'),
                         {'blurRad': '90000', 'dist': '38000', 'dir': '5400000', 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (color[0], color[1], color[2])})
    a = clr.makeelement(qn('a:alpha'), {'val': str(alpha * 1000)})
    clr.append(a); sh.append(clr); eff.append(sh); sp.append(eff)


def rect(s, x, y, w, h, fill=CARD, line=None, lw=1.0, radius=True, shadow=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if shadow:
        _soft_shadow(shp)
    return shp


def circle(s, x, y, d, fill, glyph='', gcolor=WHITE, gsize=16):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = fill; c.line.fill.background()
    c.shadow.inherit = False
    if glyph:
        tf = c.text_frame; tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = glyph
        r.font.size = Pt(gsize); r.font.color.rgb = gcolor; r.font.name = BODY
        c.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return c


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, font, italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        for (txt, size, color, bold, font, *rest) in para:
            italic = rest[0] if rest else False
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font; r.font.italic = italic
    return tb


def R(t, sz, c, b=False, f=BODY, it=False):
    return (t, sz, c, b, f, it)


def header(s, kicker, title, dark=False):
    tc = WHITE if dark else NAVY
    kc = AMBER
    text(s, 0.7, 0.5, 12, 0.4, [[R(kicker.upper(), 13, kc, True, BODY)]], space_after=0)
    text(s, 0.68, 0.82, 12, 1.0, [[R(title, 30, tc, True, HEAD)]], space_after=0)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = slide(INK)
# soft decorative circles
for (cx, cy, d, col, al) in [(11.4, -1.1, 4.2, NAVY, 1), (12.2, 4.9, 3.6, NAVY, 1), (-1.2, 5.2, 3.4, NAVY, 1)]:
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = col; o.line.fill.background(); o.shadow.inherit = False
circle(s, 0.7, 0.95, 0.95, AMBER, '📅', INK, 34)
text(s, 1.85, 1.0, 8, 0.9, [[R('EduScheduler', 30, WHITE, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.7, 2.55, 12, 1.6,
     [[R('Smart College Timetable', 50, WHITE, True, HEAD)],
      [R('Management System', 50, AMBER, True, HEAD)]], space_after=2)
text(s, 0.72, 4.55, 11.6, 0.8,
     [[R('Upload your data, click once, and get a clash-free timetable for every', 17, ICE, False, BODY)],
      [R('department, professor, room and student — automatically.', 17, ICE, False, BODY)]], space_after=2)
# pills
pills = ['Automatic Scheduling', 'Role-Based Access', 'One-Click Export']
px = 0.72
for p in pills:
    w = 0.32 + 0.108 * len(p)
    rect(s, px, 5.95, w, 0.5, fill=NAVY, radius=True)
    text(s, px, 5.95, w, 0.5, [[R(p, 12.5, ICE, True, BODY)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    px += w + 0.25
s.notes_slide.notes_text_frame.text = (
    "Welcome. EduScheduler turns the painful, manual job of building college timetables into a one-click task. "
    "Today I'll show what it does, who uses it, and why it matters — no technical background needed.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE CHALLENGE
# ════════════════════════════════════════════════════════════════════════════
s = slide(LIGHT)
header(s, 'The problem', 'Building timetables by hand is hard')
cards = [
    ('🕒', 'Days of effort', 'Coordinators spend days juggling spreadsheets every semester — and redo it all whenever one change happens.'),
    ('⚠️', 'Costly clashes', 'A professor in two rooms at once, or two classes in one room — mistakes that are easy to make and hard to spot.'),
    ('🧩', 'Too many rules', 'Lunch breaks, lab pairings, workload limits, shared faculty across departments — all must line up at once.'),
]
cw, gap = 3.93, 0.32
x0 = 0.7
for i, (ic, t, d) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rect(s, x, 2.0, cw, 3.7, fill=CARD, shadow=True)
    circle(s, x + 0.35, 2.35, 0.95, ICE, ic, NAVY, 30)
    text(s, x + 0.35, 3.55, cw - 0.7, 0.6, [[R(t, 19, NAVY, True, HEAD)]])
    text(s, x + 0.35, 4.15, cw - 0.7, 1.4, [[R(d, 13.5, SLATE, False, BODY)]], space_after=0)
text(s, 0.7, 6.05, 12, 0.7,
     [[R('The result: ', 14, INKTX, True, BODY), R('wasted time, frustrated staff, and timetables that still contain mistakes.', 14, SLATE, False, BODY, True)]],
     anchor=MSO_ANCHOR.MIDDLE)
s.notes_slide.notes_text_frame.text = (
    "Anyone who has built a timetable knows the pain: days of spreadsheet work, and a single change means starting over. "
    "Clashes are easy to create and hard to catch. EduScheduler removes that burden.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — MEET EDUSCHEDULER (3-step pipeline)
# ════════════════════════════════════════════════════════════════════════════
s = slide(WHITE)
header(s, 'The solution', 'One system, three simple steps')
steps = [
    ('📥', '1 · Upload', 'Fill in five simple spreadsheets and upload them — departments, professors, rooms, sections and subjects.'),
    ('⚡', '2 · Generate', 'Click once. The system builds a complete, clash-free timetable that respects every rule.'),
    ('👁️', '3 · View & Share', 'Everyone sees their own view — and exports to PDF, spreadsheet or a scannable QR code.'),
]
cw = 3.93
x0 = 0.7
for i, (ic, t, d) in enumerate(steps):
    x = x0 + i * (cw + 0.32)
    rect(s, x, 2.05, cw, 3.5, fill=TINT, shadow=True)
    circle(s, x + cw/2 - 0.5, 2.4, 1.0, NAVY, ic, WHITE, 32)
    text(s, x, 3.65, cw, 0.5, [[R(t, 20, NAVY, True, HEAD)]], align=PP_ALIGN.CENTER)
    text(s, x + 0.4, 4.25, cw - 0.8, 1.2, [[R(d, 13.5, SLATE, False, BODY)]], align=PP_ALIGN.CENTER, space_after=0)
    if i < 2:
        text(s, x + cw - 0.05, 3.0, 0.5, 0.8, [[R('→', 26, AMBER, True, BODY)]], align=PP_ALIGN.CENTER)
text(s, 0.7, 6.0, 12, 0.8,
     [[R('No formulas. No manual checking. The computer does the heavy lifting — you stay in control.', 15, NAVY, True, BODY)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
s.notes_slide.notes_text_frame.text = (
    "The whole system is three steps: upload your data, click generate, and view or share. "
    "That's it — the complexity is hidden inside.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — FOUR TYPES OF USERS
# ════════════════════════════════════════════════════════════════════════════
s = slide(LIGHT)
header(s, 'Who uses it', 'Four kinds of users, one login screen')
roles = [
    ('🛡️', 'Administrator', 'Full control', 'Manages every department, uploads data, creates accounts and generates timetables.', NAVY),
    ('🏛️', 'Department Admin', 'Their department only', 'Uploads and manages their own department — and cannot see or change any other.', RGBColor(0x2B,0x5C,0xA8)),
    ('👨‍🏫', 'Professor', 'Own timetable', 'Logs in with their name and ID to view only their personal teaching schedule.', RGBColor(0x3A,0x46,0x73)),
    ('🎓', 'Student', 'No password', 'Simply picks their department, course, semester and section to see their timetable.', RGBColor(0x55,0x5F,0x86)),
]
cw = 2.93
x0 = 0.7
for i, (ic, t, tag, d, col) in enumerate(roles):
    x = x0 + i * (cw + 0.29)
    rect(s, x, 2.0, cw, 4.0, fill=CARD, shadow=True)
    circle(s, x + 0.32, 2.3, 0.92, col, ic, WHITE, 26)
    text(s, x + 0.32, 3.42, cw - 0.6, 0.5, [[R(t, 16.5, NAVY, True, HEAD)]])
    rect(s, x + 0.32, 3.95, cw - 0.64, 0.4, fill=ICE, radius=True)
    text(s, x + 0.32, 3.95, cw - 0.64, 0.4, [[R(tag, 11, NAVY, True, BODY)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.32, 4.5, cw - 0.6, 1.4, [[R(d, 12.5, SLATE, False, BODY)]], space_after=0)
text(s, 0.7, 6.35, 12, 0.6,
     [[R('Everyone sees exactly what they need — and nothing they shouldn’t.', 14, INKTX, True, BODY, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
s.notes_slide.notes_text_frame.text = (
    "Four roles. Admin runs everything. Department admins are locked to their own department. "
    "Professors see only their own schedule. Students need no account at all — they just pick their class.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS (workflow)
# ════════════════════════════════════════════════════════════════════════════
s = slide(WHITE)
header(s, 'The workflow', 'From spreadsheet to finished timetable')
flow = [
    ('📝', 'Prepare', 'Download the ready-made templates and fill in your details.'),
    ('📤', 'Upload', 'Drag the five files in. The system reads and checks them.'),
    ('⚙️', 'Generate', 'One click builds the full timetable using all the rules.'),
    ('📲', 'Distribute', 'Share via dashboard, PDF, spreadsheet or QR code.'),
]
bw = 2.85
x0 = 0.72
y = 2.35
for i, (ic, t, d) in enumerate(flow):
    x = x0 + i * (bw + 0.18)
    rect(s, x, y, bw, 2.7, fill=TINT, shadow=True)
    circle(s, x + 0.3, y + 0.32, 0.86, AMBER, ic, INK, 26)
    text(s, x + 1.28, y + 0.3, bw - 1.4, 0.9, [[R(t, 17, NAVY, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.32, y + 1.35, bw - 0.6, 1.2, [[R(d, 12.5, SLATE, False, BODY)]], space_after=0)
    if i < 3:
        text(s, x + bw - 0.02, y + 0.85, 0.42, 0.8, [[R('→', 24, NAVY, True, BODY)]], align=PP_ALIGN.CENTER)
# bottom strip
rect(s, 0.72, 5.5, 11.9, 1.15, fill=NAVY, shadow=True)
text(s, 1.05, 5.5, 11.3, 1.15,
     [[R('Re-upload anytime. ', 15, AMBER, True, BODY), R('Change a professor or add a section, upload again, and the timetable rebuilds itself — other departments stay untouched.', 15, WHITE, False, BODY)]],
     anchor=MSO_ANCHOR.MIDDLE, space_after=0)
s.notes_slide.notes_text_frame.text = (
    "Prepare, upload, generate, distribute. And it's repeatable — re-uploading rebuilds the timetable instantly, "
    "without disturbing departments that didn't change.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SMART CSV IMPORT
# ════════════════════════════════════════════════════════════════════════════
s = slide(LIGHT)
header(s, 'Step 1 · Data', 'Five simple files describe your college')
files = [
    ('🏢', 'Departments', 'Working hours, start time and lunch break for each department.'),
    ('👩‍🏫', 'Professors', 'Names, IDs, workload limits, subjects, plus must-teach and unavailable times.'),
    ('🚪', 'Rooms & Labs', 'Room names, capacities and which subjects each lab can host.'),
    ('👥', 'Sections', 'Each class group, its students, fixed room and any weekly holiday.'),
    ('📚', 'Subjects', 'Theory, lab, tutorial, NPTEL and elective hours per week.'),
]
# left: file cards (2 col)
cw = 5.75
positions = [(0.7,2.0),(6.85,2.0),(0.7,3.35),(6.85,3.35),(0.7,4.7)]
for (ic,t,d),(x,yy) in zip(files, positions):
    rect(s, x, yy, cw, 1.2, fill=CARD, shadow=True)
    circle(s, x+0.25, yy+0.27, 0.68, NAVY, ic, WHITE, 20)
    text(s, x+1.1, yy+0.16, cw-1.2, 0.4, [[R(t, 15.5, NAVY, True, HEAD)]])
    text(s, x+1.1, yy+0.58, cw-1.25, 0.6, [[R(d, 11.5, SLATE, False, BODY)]], space_after=0)
# bottom-right highlight card
rect(s, 6.85, 4.7, 5.75, 1.2, fill=NAVY, shadow=True)
circle(s, 6.85+0.25, 4.7+0.27, 0.68, AMBER, '⬇️', INK, 20)
text(s, 6.85+1.1, 4.7+0.14, 5.75-1.2, 0.45, [[R('Download-ready templates', 15, WHITE, True, HEAD)]])
text(s, 6.85+1.1, 4.7+0.56, 5.75-1.25, 0.6, [[R('Pre-filled examples you just edit — no guesswork about the format.', 11.5, ICE, False, BODY)]], space_after=0)
text(s, 0.7, 6.2, 12, 0.6, [[R('Names can be messy — “3rd” or “Semester 3”, “MB-202” or “Room 202” — the system understands them.', 13.5, INKTX, True, BODY, True)]], anchor=MSO_ANCHOR.MIDDLE)
s.notes_slide.notes_text_frame.text = (
    "Five spreadsheets describe the whole college. Each has a downloadable template with examples. "
    "The importer is forgiving about formats — it understands ordinals, room codes and different spellings.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ONE-CLICK GENERATION
# ════════════════════════════════════════════════════════════════════════════
s = slide(INK)
# decorative
o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(-1.0), Inches(4.0), Inches(4.0))
o.fill.solid(); o.fill.fore_color.rgb = NAVY; o.line.fill.background(); o.shadow.inherit=False
text(s, 0.7, 0.5, 12, 0.4, [[R('STEP 2 · THE ENGINE', 13, AMBER, True, BODY)]], space_after=0)
text(s, 0.68, 0.86, 12, 1.0, [[R('One click builds the entire timetable', 30, WHITE, True, HEAD)]], space_after=0)
circle(s, 0.7, 2.2, 1.5, AMBER, '⚡', INK, 52)
text(s, 2.55, 2.25, 10, 1.5,
     [[R('Behind the button is a scheduling engine that places every class for every', 16, ICE, False, BODY)],
      [R('group — while checking dozens of rules at the same time, in seconds.', 16, ICE, False, BODY)]], anchor=MSO_ANCHOR.MIDDLE, space_after=3)
# three mini stats
stats = [('Every', 'class placed automatically'), ('Zero', 'double-bookings'), ('Seconds', 'not days of work')]
cw = 3.93
for i,(big,small) in enumerate(stats):
    x = 0.7 + i*(cw+0.32)
    rect(s, x, 4.2, cw, 1.9, fill=NAVY, shadow=True)
    text(s, x, 4.45, cw, 0.9, [[R(big, 34, AMBER, True, HEAD)]], align=PP_ALIGN.CENTER)
    text(s, x+0.3, 5.4, cw-0.6, 0.6, [[R(small, 14, WHITE, False, BODY)]], align=PP_ALIGN.CENTER, space_after=0)
text(s, 0.7, 6.55, 12, 0.5, [[R('If a class genuinely cannot fit, you get a clear warning — never a silent mistake.', 13.5, ICE, False, BODY, True)]], align=PP_ALIGN.CENTER)
s.notes_slide.notes_text_frame.text = (
    "One button. The engine places every class while juggling all the rules at once, in seconds. "
    "And if something truly can't fit, it tells you — it never hides a problem.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — THE SMART RULES
# ════════════════════════════════════════════════════════════════════════════
s = slide(WHITE)
header(s, 'The intelligence', 'The rules it follows — automatically')
rules = [
    ('🚫', 'No back-to-back overload', 'Never three classes — or three labs — in a row for a professor.'),
    ('⚖️', 'Fair workload', 'Teaching is spread evenly across the week and across staff.'),
    ('🍽️', 'Lunch respected', 'No classes during each department’s lunch break; online (NPTEL) classes only after lunch.'),
    ('🔗', 'No clashes anywhere', 'One room, one class. A shared professor is never in two places at once.'),
    ('📌', 'Fixed & blocked times', 'Honours “must teach at this time” and “unavailable” slots from the data.'),
    ('🏫', 'Smart room choice', 'Uses a section’s fixed room, or finds a free room in the same department.'),
]
cw, ch = 3.93, 1.95
gx, gy = 0.32, 0.28
for i,(ic,t,d) in enumerate(rules):
    col = i % 3; row = i // 3
    x = 0.7 + col*(cw+gx); y = 2.0 + row*(ch+gy)
    rect(s, x, y, cw, ch, fill=TINT, shadow=True)
    circle(s, x+0.28, y+0.3, 0.72, NAVY, ic, WHITE, 21)
    text(s, x+1.12, y+0.28, cw-1.2, 0.7, [[R(t, 14.5, NAVY, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+0.3, y+1.05, cw-0.55, 0.85, [[R(d, 11.5, SLATE, False, BODY)]], space_after=0)
s.notes_slide.notes_text_frame.text = (
    "These are the rules a human would try to remember. The system enforces all of them at once: "
    "no overload, fair workload, lunch breaks, no clashes, fixed and blocked times, and smart room choice.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — LABS, TUTORIALS & ELECTIVES
# ════════════════════════════════════════════════════════════════════════════
s = slide(LIGHT)
header(s, 'Special cases', 'Labs, tutorials and electives — handled')
blocks = [
    ('🧪', 'Two-hour labs', 'A lab takes two back-to-back periods in a proper lab room — never split up.'),
    ('📝', 'Parallel tutorials', 'While one group is in the lab, other groups run their tutorials at the same time.'),
    ('🔀', 'Electives together', 'A semester’s electives run in the same slot, in different rooms — so a student simply picks one.'),
]
cw = 3.93
for i,(ic,t,d) in enumerate(blocks):
    x = 0.7 + i*(cw+0.32)
    rect(s, x, 2.0, cw, 2.6, fill=CARD, shadow=True)
    circle(s, x+0.35, 2.3, 0.95, ICE, ic, NAVY, 28)
    text(s, x+1.45, 2.35, cw-1.5, 0.85, [[R(t, 16, NAVY, True, HEAD)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+0.36, 3.45, cw-0.7, 1.0, [[R(d, 12.5, SLATE, False, BODY)]], space_after=0)
# illustrative timetable strip
rect(s, 0.7, 5.0, 11.93, 1.65, fill=WHITE, shadow=True)
text(s, 0.95, 5.1, 11, 0.35, [[R('Example — one two-hour block, three groups all busy:', 12.5, NAVY, True, BODY)]])
cells = [('Group A','Lab  (Room L1)',GREEN),('Group B','Tutorial (R2)',NAVY),('Group C','Tutorial (R3)',RGBColor(0x7C,0x3A,0xED))]
bx = 0.95
for (g, lab, col) in cells:
    rect(s, bx, 5.55, 3.7, 0.95, fill=col, radius=True)
    text(s, bx, 5.62, 3.7, 0.4, [[R(g, 12.5, WHITE, True, BODY)]], align=PP_ALIGN.CENTER)
    text(s, bx, 6.0, 3.7, 0.4, [[R(lab, 11.5, WHITE, False, BODY)]], align=PP_ALIGN.CENTER)
    bx += 3.9
s.notes_slide.notes_text_frame.text = (
    "Labs run as proper two-hour blocks, and while one group is in the lab the others do tutorials at the same time, "
    "so nobody is idle. Electives for a semester run together so students just choose one.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CROSS-DEPARTMENT
# ════════════════════════════════════════════════════════════════════════════
s = slide(WHITE)
header(s, 'Coordination', 'Departments that don’t collide')
text(s, 0.7, 1.95, 7.2, 2.2,
     [[R('Colleges share professors and rooms across departments. EduScheduler treats the whole college as one picture, so:', 15, INKTX, False, BODY)]], space_after=6)
points = [
    ('A room is never double-booked, even across two departments.'),
    ('A professor teaching in two departments is never scheduled twice at once.'),
    ('When one department re-uploads, the shared professor’s timetable updates — and the others are left intact.'),
]
yy = 3.05
for p in points:
    circle(s, 0.75, yy+0.02, 0.36, GREEN, '✓', WHITE, 14)
    text(s, 1.3, yy-0.05, 6.7, 0.8, [[R(p, 13.5, INKTX, False, BODY)]], space_after=0)
    yy += 0.95
# right visual
rect(s, 8.35, 1.95, 4.28, 4.2, fill=TINT, shadow=True)
text(s, 8.6, 2.15, 3.8, 0.4, [[R('Shared resources', 14, NAVY, True, HEAD)]])
shared = [('👤  Dr. Sharma','teaches CSE + ECE',NAVY),('🚪  Seminar Hall','used by 3 departments',RGBColor(0x2B,0x5C,0xA8))]
ry = 2.7
for (t, d, col) in shared:
    rect(s, 8.6, ry, 3.78, 0.95, fill=CARD)
    text(s, 8.8, ry+0.13, 3.5, 0.4, [[R(t, 13.5, col, True, BODY)]])
    text(s, 8.8, ry+0.5, 3.5, 0.35, [[R(d, 11.5, SLATE, False, BODY)]])
    ry += 1.1
rect(s, 8.6, ry, 3.78, 0.95, fill=GREEN)
text(s, 8.6, ry, 3.78, 0.95, [[R('✓  Always conflict-free', 13.5, WHITE, True, BODY)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
s.notes_slide.notes_text_frame.text = (
    "Real colleges share staff and rooms. The system sees the whole college at once, so shared rooms and shared "
    "professors are never double-booked — and one department's upload never breaks another's timetable.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — VIEW & EXPORT
# ════════════════════════════════════════════════════════════════════════════
s = slide(LIGHT)
header(s, 'Step 3 · Share', 'Every view, every format')
views = [
    ('🗂️', 'By section', 'Each class group’s weekly grid.'),
    ('👨‍🏫', 'By professor', 'A teacher’s personal schedule with fixed/blocked times marked.'),
    ('🚪', 'By room', 'What happens in each room, all week.'),
    ('🎓', 'By semester', 'A whole semester at a glance.'),
]
cw = 2.93
for i,(ic,t,d) in enumerate(views):
    x = 0.7 + i*(cw+0.29)
    rect(s, x, 2.0, cw, 2.35, fill=CARD, shadow=True)
    circle(s, x+cw/2-0.45, 2.25, 0.9, NAVY, ic, WHITE, 26)
    text(s, x, 3.3, cw, 0.4, [[R(t, 15, NAVY, True, HEAD)]], align=PP_ALIGN.CENTER)
    text(s, x+0.3, 3.75, cw-0.6, 0.6, [[R(d, 11.5, SLATE, False, BODY)]], align=PP_ALIGN.CENTER, space_after=0)
# export row
rect(s, 0.7, 4.75, 11.93, 1.7, fill=NAVY, shadow=True)
text(s, 1.0, 4.95, 11, 0.45, [[R('Export & share in one click', 17, WHITE, True, HEAD)]])
exps = [('📄','PDF','Print-ready timetables'),('📊','Spreadsheet','Open in Excel'),('📱','QR code','Scan to view on a phone')]
ex = 1.0
for (ic, t, d) in exps:
    circle(s, ex, 5.55, 0.72, AMBER, ic, INK, 22)
    text(s, ex+0.85, 5.5, 2.9, 0.4, [[R(t, 14, WHITE, True, BODY)]])
    text(s, ex+0.85, 5.9, 2.9, 0.4, [[R(d, 11, ICE, False, BODY)]])
    ex += 3.95
s.notes_slide.notes_text_frame.text = (
    "Everyone gets the view they need — by section, professor, room or semester — and any of them exports to "
    "PDF, spreadsheet, or a QR code students can scan with a phone.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SECURITY / ACCESS TABLE
# ════════════════════════════════════════════════════════════════════════════
s = slide(WHITE)
header(s, 'Safety', 'Everyone sees only what they should')
rows = [
    ('Capability', 'Admin', 'Dept Admin', 'Professor', 'Student'),
    ('Upload data & generate', '✓ all', '✓ own dept', '—', '—'),
    ('Manage accounts', '✓', '—', '—', '—'),
    ('View any timetable', '✓', 'own dept', 'own only', 'own only'),
    ('Needs a password', '✓', '✓', '✓ (ID)', 'No'),
]
tx, ty = 0.7, 2.15
colw = [4.6, 1.83, 1.95, 1.75, 1.7]
rh = 0.78
for ri, row in enumerate(rows):
    cx = tx
    for ci, val in enumerate(row):
        head = (ri == 0)
        fill = NAVY if head else (TINT if ri % 2 else WHITE)
        rect(s, cx, ty + ri*rh, colw[ci]-0.06, rh-0.08, fill=fill, radius=True)
        col = WHITE if head else (NAVY if ci == 0 else INKTX)
        bold = head or ci == 0
        sz = 13 if head else 12.5
        al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        padx = 0.2 if ci == 0 else 0
        text(s, cx+padx, ty + ri*rh, colw[ci]-0.06-padx, rh-0.08, [[R(str(val), sz, col, bold, BODY)]],
             align=al, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        cx += colw[ci]
text(s, 0.7, 6.55, 12, 0.5, [[R('A department admin literally cannot open another department’s pages — access is enforced, not just hidden.', 13, SLATE, False, BODY, True)]], align=PP_ALIGN.CENTER)
s.notes_slide.notes_text_frame.text = (
    "Access is enforced at every page. A department admin can't even open another department's pages. "
    "Students need no password; professors log in with their name and ID.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — BENEFITS (dark, stat callouts)
# ════════════════════════════════════════════════════════════════════════════
s = slide(INK)
text(s, 0.7, 0.5, 12, 0.4, [[R('WHY IT MATTERS', 13, AMBER, True, BODY)]], space_after=0)
text(s, 0.68, 0.86, 12, 1.0, [[R('Days of work become a single click', 30, WHITE, True, HEAD)]], space_after=0)
big = [
    ('Days → seconds', 'Timetables that took days are generated in moments.'),
    ('0 clashes', 'No double-booked rooms or professors — guaranteed by design.'),
    ('4 roles', 'One secure system for admins, departments, staff and students.'),
    ('1 click', 'Re-generate the whole college whenever anything changes.'),
]
cw = 5.85
for i,(b,d) in enumerate(big):
    col = i % 2; row = i // 2
    x = 0.7 + col*(cw+0.33); y = 2.25 + row*2.05
    rect(s, x, y, cw, 1.8, fill=NAVY, shadow=True)
    text(s, x+0.4, y+0.22, cw-0.8, 0.85, [[R(b, 32, AMBER, True, HEAD)]])
    text(s, x+0.42, y+1.05, cw-0.85, 0.6, [[R(d, 14, ICE, False, BODY)]], space_after=0)
s.notes_slide.notes_text_frame.text = (
    "The bottom line: days of work become a click, clashes are designed out, four roles are served by one secure "
    "system, and the whole college can be regenerated instantly.")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CLOSING
# ════════════════════════════════════════════════════════════════════════════
s = slide(INK)
for (cx, cy, d) in [(11.2, -1.2, 4.4), (-1.4, 4.8, 3.8)]:
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = NAVY; o.line.fill.background(); o.shadow.inherit=False
circle(s, 0.7, 1.3, 1.0, AMBER, '📅', INK, 34)
text(s, 0.7, 2.7, 12, 1.2, [[R('Smarter timetables,', 46, WHITE, True, HEAD)],[R('zero headaches.', 46, AMBER, True, HEAD)]], space_after=2)
text(s, 0.72, 4.85, 11.5, 0.8, [[R('Upload once. Generate in a click. Share with everyone.', 18, ICE, False, BODY)]])
rect(s, 0.72, 5.85, 4.9, 0.75, fill=AMBER, radius=True)
text(s, 0.72, 5.85, 4.9, 0.75, [[R('Thank you  ·  Questions?', 16, INK, True, BODY)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
s.notes_slide.notes_text_frame.text = (
    "To wrap up: EduScheduler makes timetabling smart and painless. Upload once, generate in a click, "
    "and everyone gets their schedule. Happy to take questions.")

out = "D:/college_timetable_v2_work/_deliverables/EduScheduler_Presentation.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
